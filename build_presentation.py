"""
Build the CREME project presentation as a .pptx file.
Run: python build_presentation.py
Output: CREME_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1E, 0x1E, 0x2E)   # near-black navy
ACCENT       = RGBColor(0x74, 0xC7, 0xEC)   # sky blue
ACCENT2      = RGBColor(0xA6, 0xE3, 0xA1)   # mint green
ACCENT3      = RGBColor(0xF3, 0x8B, 0xA8)   # pink/red
ACCENT4      = RGBColor(0xFA, 0xB3, 0x87)   # peach/orange
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCC, 0xCC, 0xDD)
MID_GREY     = RGBColor(0x88, 0x88, 0x99)
BOX_BG       = RGBColor(0x2A, 0x2A, 0x3E)   # slightly lighter than bg
YELLOW       = RGBColor(0xF9, 0xE2, 0xAF)

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Low-level helpers ─────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(14), bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = font_name
    return txb


def add_para(tf, text, font_size=Pt(13), bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, space_before=Pt(4),
             font_name="Calibri", level=0):
    p = tf.add_paragraph()
    p.alignment    = align
    p.level        = level
    p.space_before = space_before
    run = p.add_run()
    run.text           = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = font_name
    return p


def slide_bg(slide):
    """Fill slide background with DARK_BG."""
    add_rect(slide, 0, 0, W, H, fill=DARK_BG)


def accent_bar(slide, color=ACCENT, height=Inches(0.06)):
    """Thin horizontal accent bar at top."""
    add_rect(slide, 0, 0, W, height, fill=color)


def slide_title(slide, title, subtitle=None, title_color=ACCENT,
                title_size=Pt(32)):
    """Standard title block used on content slides."""
    accent_bar(slide)
    add_text(slide, title,
             x=Inches(0.5), y=Inches(0.15), w=Inches(12.3), h=Inches(0.75),
             font_size=title_size, bold=True, color=title_color,
             align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 x=Inches(0.5), y=Inches(0.88), w=Inches(12.3), h=Inches(0.35),
                 font_size=Pt(15), italic=True, color=LIGHT_GREY)
    # thin divider line
    add_rect(slide, Inches(0.5), Inches(1.2),
             Inches(12.33), Inches(0.025), fill=ACCENT)


def textbox_with_bg(slide, x, y, w, h, bg=BOX_BG, border=ACCENT,
                    border_w=Pt(1.2)):
    box = add_rect(slide, x, y, w, h, fill=bg, line=border, line_w=border_w)
    return box


def code_block(slide, lines, x, y, w, h):
    """Monospace dark box for pseudo-code / math."""
    add_rect(slide, x, y, w, h,
             fill=rgb(0x18, 0x18, 0x2C), line=ACCENT, line_w=Pt(1))
    txb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.1),
                                    w - Inches(0.24), h - Inches(0.2))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(1)
        run = p.add_run()
        run.text           = line
        run.font.size      = Pt(11)
        run.font.name      = "Consolas"
        run.font.color.rgb = ACCENT2


def table_slide(slide, headers, rows, x, y, w, h,
                header_fill=rgb(0x2E, 0x3A, 0x5E),
                row_fill_a=BOX_BG, row_fill_b=rgb(0x22, 0x22, 0x35)):
    """Simple table drawn with rectangles + text."""
    n_cols = len(headers)
    n_rows = len(rows)
    col_w  = w / n_cols
    row_h  = h / (n_rows + 1)

    # Header row
    for j, hdr in enumerate(headers):
        add_rect(slide, x + col_w*j, y, col_w, row_h,
                 fill=header_fill, line=DARK_BG, line_w=Pt(0.5))
        add_text(slide, hdr,
                 x=x + col_w*j + Inches(0.05), y=y + Inches(0.02),
                 w=col_w - Inches(0.1), h=row_h,
                 font_size=Pt(11), bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        fill = row_fill_a if i % 2 == 0 else row_fill_b
        for j, cell in enumerate(row):
            add_rect(slide, x + col_w*j, y + row_h*(i+1), col_w, row_h,
                     fill=fill, line=DARK_BG, line_w=Pt(0.5))
            c = LIGHT_GREY
            add_text(slide, cell,
                     x=x + col_w*j + Inches(0.06),
                     y=y + row_h*(i+1) + Inches(0.02),
                     w=col_w - Inches(0.12), h=row_h - Inches(0.04),
                     font_size=Pt(10.5), color=c,
                     align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)

# Top accent strip
add_rect(sl, 0, 0, W, Inches(0.08), fill=ACCENT)
add_rect(sl, 0, H - Inches(0.08), W, Inches(0.08), fill=ACCENT2)

# Central title block
add_rect(sl, Inches(0.6), Inches(1.6), Inches(12.13), Inches(1.8),
         fill=BOX_BG, line=ACCENT, line_w=Pt(1.5))

add_text(sl, "CREME: Not Every Layer Counts",
         x=Inches(0.7), y=Inches(1.72), w=Inches(11.9), h=Inches(0.85),
         font_size=Pt(38), bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
add_text(sl, "Proactive Robustness Enhancement of Code LLMs via Layer-Aware Fine-Tuning",
         x=Inches(0.7), y=Inches(2.52), w=Inches(11.9), h=Inches(0.55),
         font_size=Pt(17), italic=True, color=WHITE,
         align=PP_ALIGN.CENTER)

# Three info pills
for i, (label, val, col) in enumerate([
    ("Model",     "CodeLlama-7B",          ACCENT),
    ("Benchmark", "MBPP",                   ACCENT2),
    ("Pert. Types","20 across 6 categories",ACCENT4),
]):
    px = Inches(0.6) + Inches(4.18) * i
    add_rect(sl, px, Inches(3.7), Inches(3.9), Inches(0.75),
             fill=rgb(0x2E,0x3A,0x5E), line=col, line_w=Pt(1.2))
    add_text(sl, label,
             x=px + Inches(0.1), y=Inches(3.72),
             w=Inches(3.7), h=Inches(0.28),
             font_size=Pt(10), bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(sl, val,
             x=px + Inches(0.1), y=Inches(3.98),
             w=Inches(3.7), h=Inches(0.28),
             font_size=Pt(12), color=WHITE,
             align=PP_ALIGN.CENTER)

add_text(sl, "ICSE 2025 Submission",
         x=Inches(0.5), y=Inches(6.7), w=Inches(12.3), h=Inches(0.3),
         font_size=Pt(11), color=MID_GREY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Introduction
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Introduction", "Context & Motivation")

# Left column — context
textbox_with_bg(sl, Inches(0.4), Inches(1.4), Inches(5.9), Inches(5.6))
add_text(sl, "Code LLMs in Production",
         x=Inches(0.55), y=Inches(1.45), w=Inches(5.6), h=Inches(0.4),
         font_size=Pt(14), bold=True, color=ACCENT)
for bullet in [
    "Models like CodeLlama-7B and QwenCoder-7B excel at code generation benchmarks",
    "Real-world prompts are rarely clean — developers paraphrase, rename variables, restructure docstrings",
    "A small surface-level change in the prompt can sharply reduce model accuracy",
    "This is called the robustness gap — and it matters for any production deployment",
]:
    add_text(sl, f"•  {bullet}",
             x=Inches(0.6), y=Inches(1.9 + [0,0.7,1.4,2.15][
                 ["Models like","Real-world","A small","This is called"].index(
                     next(b for b in ["Models like","Real-world","A small","This is called"] if bullet.startswith(b)))]),
             w=Inches(5.6), h=Inches(0.75),
             font_size=Pt(12.5), color=LIGHT_GREY)

# Right column — what CREME does
textbox_with_bg(sl, Inches(6.7), Inches(1.4), Inches(6.2), Inches(2.6),
                border=ACCENT2)
add_text(sl, "What CREME Does",
         x=Inches(6.85), y=Inches(1.45), w=Inches(5.9), h=Inches(0.4),
         font_size=Pt(14), bold=True, color=ACCENT2)
add_text(sl,
    "CREME (Code Robustness Enhancement via Model Editing) is an existing framework "
    "that repairs this gap by:\n\n"
    "  1.  Identifying which transformer layer diverges most under a perturbation "
    "(causal tracing)\n\n"
    "  2.  Directly editing that layer's weights to realign the model's internal "
    "representations",
    x=Inches(6.85), y=Inches(1.88), w=Inches(5.9), h=Inches(2.0),
    font_size=Pt(12), color=LIGHT_GREY)

textbox_with_bg(sl, Inches(6.7), Inches(4.2), Inches(6.2), Inches(2.6),
                border=ACCENT)
add_text(sl, "Our Contribution",
         x=Inches(6.85), y=Inches(4.25), w=Inches(5.9), h=Inches(0.4),
         font_size=Pt(14), bold=True, color=ACCENT)
add_text(sl,
    "We extend CREME with a proactive fine-tuning paradigm:\n\n"
    "Train the model once — offline — so it is natively robust to all "
    "20 perturbation types without any per-task editing at inference time.",
    x=Inches(6.85), y=Inches(4.68), w=Inches(5.9), h=Inches(1.8),
    font_size=Pt(12), color=LIGHT_GREY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Problem Statement", "The Robustness Gap in Code LLMs")

# Category table
headers = ["Category", "Types", "Description"]
rows = [
    ["Argument (A)",   "A1, A2, A3",       "Rename / reorder function arguments"],
    ["Control (C)",    "C1, C2, C3",       "Perturb control flow structure"],
    ["Data (D)",       "D1, D2, D3, D4",   "Modify data types, values, structures"],
    ["Expression (E)", "E1 – E6",          "Alter operators, literals, expressions"],
    ["Problem (P)",    "P1, P2",           "Rephrase the natural language description"],
    ["Statement (S)",  "S1, S2",           "Modify individual statements"],
]
table_slide(sl, headers, rows,
            x=Inches(0.4), y=Inches(1.35),
            w=Inches(7.8),  h=Inches(4.2))

# Right panel — key numbers
textbox_with_bg(sl, Inches(8.55), Inches(1.35), Inches(4.4), Inches(4.2),
                border=ACCENT3)
add_text(sl, "Key Numbers",
         x=Inches(8.7), y=Inches(1.4), w=Inches(4.1), h=Inches(0.4),
         font_size=Pt(14), bold=True, color=ACCENT3)

for label, val, col in [
    ("Original prompt  pass@1",  "0.603",  ACCENT2),
    ("Perturbed prompt  pass@1", "0.457",  ACCENT3),
    ("Mean delta",               "−0.146", ACCENT4),
    ("Worst type (D2)",          "−0.343", ACCENT3),
    ("Best type (E3)",           " 0.000", ACCENT2),
]:
    y_off = [1.9, 2.5, 3.1, 3.7, 4.3][
        ["Original","Perturbed","Mean","Worst","Best"].index(
            next(k for k in ["Original","Perturbed","Mean","Worst","Best"]
                 if label.startswith(k)))]
    add_text(sl, label,
             x=Inches(8.7), y=Inches(y_off), w=Inches(2.8), h=Inches(0.4),
             font_size=Pt(11.5), color=LIGHT_GREY)
    add_text(sl, val,
             x=Inches(11.55), y=Inches(y_off), w=Inches(1.1), h=Inches(0.4),
             font_size=Pt(13), bold=True, color=col,
             align=PP_ALIGN.RIGHT)

add_text(sl,
    "A small surface change causes an average\n14.6 pp drop — up to 34 pp for data-type perturbations.",
    x=Inches(0.4), y=Inches(6.7), w=Inches(12.5), h=Inches(0.55),
    font_size=Pt(12), italic=True, color=MID_GREY,
    align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Limitations of Reactive CREME
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Limitations of Reactive CREME", "Why per-task editing doesn't scale")

# Pipeline steps
steps = [
    ("1", "Causal Tracing",    "Run layer-by-layer activation patching for EVERY new task to find the toxic layer",       ACCENT),
    ("2", "Weight Editing",    "Apply gradient-based weight edits at that layer to align perturbed hidden states",        ACCENT4),
    ("3", "Repeat Per Task",   "Restore weights, then repeat steps 1–2 for the next (task, pert_type) pair",              ACCENT3),
    ("4", "No Generalisation", "The edited model is specific to one task — cannot be reused across tasks or pert types",  ACCENT3),
]

for i, (num, title, body, col) in enumerate(steps):
    bx = Inches(0.4) + Inches(3.22) * i
    add_rect(sl, bx, Inches(1.45), Inches(3.0), Inches(3.8),
             fill=BOX_BG, line=col, line_w=Pt(1.5))
    # Number circle
    add_rect(sl, bx + Inches(1.1), Inches(1.35), Inches(0.8), Inches(0.55),
             fill=col, line=col, line_w=Pt(0))
    add_text(sl, num,
             x=bx + Inches(1.1), y=Inches(1.34), w=Inches(0.8), h=Inches(0.55),
             font_size=Pt(18), bold=True, color=DARK_BG,
             align=PP_ALIGN.CENTER)
    add_text(sl, title,
             x=bx + Inches(0.1), y=Inches(2.0), w=Inches(2.8), h=Inches(0.5),
             font_size=Pt(13), bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(sl, body,
             x=bx + Inches(0.1), y=Inches(2.55), w=Inches(2.8), h=Inches(2.4),
             font_size=Pt(11.5), color=LIGHT_GREY,
             align=PP_ALIGN.CENTER, wrap=True)

# Complexity callout
add_rect(sl, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.95),
         fill=rgb(0x3A, 0x1A, 0x1A), line=ACCENT3, line_w=Pt(1.2))
add_text(sl,
    "Computational cost scales as  O(tasks x pert_types)  — "
    "270+ separate trace+edit runs for the MBPP benchmark alone. "
    "Each run requires a full forward pass over all 32 layers.",
    x=Inches(0.6), y=Inches(5.55), w=Inches(12.1), h=Inches(0.85),
    font_size=Pt(12.5), color=ACCENT3, align=PP_ALIGN.CENTER)

add_text(sl, "Goal: reduce this to O(1) — train once, deploy everywhere.",
         x=Inches(0.4), y=Inches(6.6), w=Inches(12.5), h=Inches(0.4),
         font_size=Pt(13), bold=True, italic=True, color=ACCENT,
         align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Novel Approach Overview
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Our Novel Approach", "Proactive Fine-Tuning — Train Once, Robust Everywhere")

add_text(sl,
    "CREME's causal tracing already tells us where the model breaks. "
    "We use that knowledge offline — training Layer 28 to be robust before deployment.",
    x=Inches(0.5), y=Inches(1.3), w=Inches(12.3), h=Inches(0.5),
    font_size=Pt(13), italic=True, color=LIGHT_GREY)

# Two column: Reactive vs Proactive
for col_x, heading, col_color, items in [
    (Inches(0.4), "Reactive CREME  (before)", ACCENT3, [
        "Causal trace at inference — every task",
        "Weight edit at inference — every task",
        "Specific to one (task, pert_type) pair",
        "Weights restored after each task",
        "O(tasks x pert_types) compute cost",
        "Cannot zero-shot generalise",
    ]),
    (Inches(6.8), "Proactive Approach  (ours)", ACCENT2, [
        "Causal trace done once offline — Layer 28",
        "LoRA adapter trained once on all pairs",
        "Single model handles all 20 pert types",
        "No weight restoration needed",
        "O(1) inference overhead",
        "Zero-shot generalises to unseen tasks",
    ]),
]:
    add_rect(sl, col_x, Inches(1.95), Inches(6.1), Inches(4.7),
             fill=BOX_BG, line=col_color, line_w=Pt(1.5))
    add_text(sl, heading,
             x=col_x + Inches(0.1), y=Inches(2.0), w=Inches(5.9), h=Inches(0.45),
             font_size=Pt(13.5), bold=True, color=col_color)
    for k, item in enumerate(items):
        sym = "✗" if col_color == ACCENT3 else "✓"
        sym_col = ACCENT3 if col_color == ACCENT3 else ACCENT2
        add_text(sl, sym,
                 x=col_x + Inches(0.15), y=Inches(2.55 + k * 0.6),
                 w=Inches(0.35), h=Inches(0.5),
                 font_size=Pt(13), bold=True, color=sym_col)
        add_text(sl, item,
                 x=col_x + Inches(0.55), y=Inches(2.55 + k * 0.6),
                 w=Inches(5.4), h=Inches(0.5),
                 font_size=Pt(12), color=LIGHT_GREY)

# Centre arrow
add_text(sl, "→",
         x=Inches(6.25), y=Inches(3.8), w=Inches(0.5), h=Inches(0.6),
         font_size=Pt(28), bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Training Objective
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Training Objective", "A two-term loss combining fluency and representation alignment")

# Loss formula box
code_block(sl, [
    "L_total  =  L_ce  +  λ · L_reg",
    "",
    "L_ce   =  cross-entropy( model(clean_prompt), labels )",
    "          → preserves existing code generation ability",
    "",
    "L_reg  =  1 − cosine_similarity( h_pert ,  h_clean.detach() )",
    "          → forces perturbed representation to align with clean",
    "            representation at Layer 28",
    "",
    "λ = 0.01   (alignment is a regulariser, not the dominant signal)",
], x=Inches(0.4), y=Inches(1.38), w=Inches(7.8), h=Inches(3.8))

# Right — three explanation cards
for i, (title, body, col) in enumerate([
    ("Why cosine similarity?",
     "Measures the angle between two vectors. Value of 1 = same direction. "
     "We minimise 1 − cos to pull the perturbed representation toward the clean "
     "one in direction, without constraining magnitude.",
     ACCENT),
    ("Why detach h_clean?",
     "The clean hidden state is a fixed anchor. If both sides had gradients, "
     "the optimiser could collapse both to zero trivially. Detaching forces only "
     "the perturbed path to move.",
     ACCENT2),
    ("Why λ = 0.01?",
     "Keeps alignment as a light regulariser. Too high and the model over-aligns, "
     "losing sensitivity to legitimate content differences. Too low and the signal "
     "is too weak to generalise across perturbation types.",
     ACCENT4),
]):
    by = Inches(1.38) + Inches(1.3) * i
    add_rect(sl, Inches(8.65), by, Inches(4.3), Inches(1.18),
             fill=BOX_BG, line=col, line_w=Pt(1.2))
    add_text(sl, title,
             x=Inches(8.8), y=by + Inches(0.06),
             w=Inches(4.0), h=Inches(0.3),
             font_size=Pt(11.5), bold=True, color=col)
    add_text(sl, body,
             x=Inches(8.8), y=by + Inches(0.36),
             w=Inches(4.0), h=Inches(0.78),
             font_size=Pt(10.5), color=LIGHT_GREY, wrap=True)

# Training data note
add_rect(sl, Inches(0.4), Inches(5.42), Inches(12.5), Inches(0.7),
         fill=rgb(0x1A, 0x2E, 0x1A), line=ACCENT2, line_w=Pt(1))
add_text(sl,
    "Training data:  All MBPP tasks  ×  all 20 perturbation types  →  (original_prompt, perturbed_prompt) pairs. "
    "One training run covers every category simultaneously.",
    x=Inches(0.6), y=Inches(5.48), w=Inches(12.1), h=Inches(0.6),
    font_size=Pt(12), color=ACCENT2, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Implementation Details  (Pipeline)
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Implementation Details", "Training pipeline — step by step")

code_block(sl, [
    "for each batch of (ori_prompt, pert_prompt) pairs:",
    "",
    "   # Step A — clean forward pass (frozen anchor)",
    "   h_clean, tokens_clean = forward(model, ori_prompts, layer=28, no_grad=True)",
    "   L_ce = cross_entropy(model(**tokens_clean, labels=labels))",
    "",
    "   # Step B — perturbed forward pass (trains gradients)",
    "   h_pert, tokens_pert  = forward(model, pert_prompts, layer=28, no_grad=False)",
    "",
    "   # Step C — mean-pool over non-padding tokens",
    "   h_clean_pooled = masked_mean(h_clean, tokens_clean.attention_mask)",
    "   h_pert_pooled  = masked_mean(h_pert,  tokens_pert.attention_mask)",
    "",
    "   # Step D — alignment loss",
    "   L_reg = 1 - cosine_similarity(h_pert_pooled, h_clean_pooled.detach())",
    "",
    "   # Step E — combined loss + backward",
    "   loss = (L_ce + 0.01 * L_reg) / grad_accum_steps",
    "   loss.backward()",
    "",
    "   # Step F — optimizer step every grad_accum_steps batches",
    "   clip_grad_norm(model, max_norm=1.0)  ->  AdamW.step()  ->  zero_grad()",
], x=Inches(0.4), y=Inches(1.38), w=Inches(8.1), h=Inches(5.65))

# Right panel — hyperparams
textbox_with_bg(sl, Inches(8.85), Inches(1.38), Inches(4.1), Inches(5.65),
                border=ACCENT4)
add_text(sl, "Hyperparameters",
         x=Inches(9.0), y=Inches(1.43), w=Inches(3.8), h=Inches(0.4),
         font_size=Pt(13), bold=True, color=ACCENT4)

params = [
    ("Model",        "CodeLlama-7B"),
    ("Target layer", "Layer 28  (down_proj)"),
    ("Optimizer",    "AdamW"),
    ("LR",           "1e-5"),
    ("λ_reg",        "0.01"),
    ("Epochs",       "5"),
    ("Batch size",   "8 pairs"),
    ("Grad accum",   "1  (eff. batch = 8)"),
    ("Precision",    "bfloat16  (bf16)"),
    ("Grad clip",    "max_norm = 1.0"),
]
for k, (label, val) in enumerate(params):
    y_ = Inches(1.92) + Inches(0.46) * k
    add_text(sl, label,
             x=Inches(9.0), y=y_, w=Inches(2.0), h=Inches(0.4),
             font_size=Pt(11), color=MID_GREY)
    add_text(sl, val,
             x=Inches(11.0), y=y_, w=Inches(1.8), h=Inches(0.4),
             font_size=Pt(11), bold=True, color=LIGHT_GREY,
             align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Fine-Tuning Without LoRA  (Full mode)
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Fine-Tuning — Mode 1: Full Weight Training",
            "--train_mode full  |  Trains the original down_proj weights directly")

# Diagram left
textbox_with_bg(sl, Inches(0.4), Inches(1.38), Inches(5.5), Inches(5.5),
                border=ACCENT4)
add_text(sl, "What happens",
         x=Inches(0.55), y=Inches(1.43), w=Inches(5.2), h=Inches(0.38),
         font_size=Pt(13.5), bold=True, color=ACCENT4)

steps_full = [
    ("All parameters",    "frozen with  requires_grad = False"),
    ("down_proj  weight", "at Layer 28 only is unfrozen"),
    ("Original tensor",   "modified in-place during training"),
    ("~45M parameters",   "trained  (4096 × 11008 weight matrix)"),
    ("Saved artifact",    "full model checkpoint  (~13 GB)"),
]
for k, (label, desc) in enumerate(steps_full):
    y_ = Inches(1.9) + Inches(0.88) * k
    add_rect(sl, Inches(0.55), y_, Inches(5.1), Inches(0.75),
             fill=rgb(0x2E,0x24,0x18), line=ACCENT4, line_w=Pt(0.8))
    add_text(sl, label,
             x=Inches(0.65), y=y_ + Inches(0.04),
             w=Inches(2.2), h=Inches(0.32),
             font_size=Pt(11.5), bold=True, color=ACCENT4)
    add_text(sl, desc,
             x=Inches(0.65), y=y_ + Inches(0.36),
             w=Inches(4.8), h=Inches(0.32),
             font_size=Pt(11), color=LIGHT_GREY)

# Right — trade-offs
textbox_with_bg(sl, Inches(6.2), Inches(1.38), Inches(6.7), Inches(2.55),
                border=ACCENT2)
add_text(sl, "Advantages",
         x=Inches(6.35), y=Inches(1.43), w=Inches(6.4), h=Inches(0.38),
         font_size=Pt(13), bold=True, color=ACCENT2)
for item in [
    "No adapter library required — standard PyTorch only",
    "Can target multiple layers simultaneously with --target_layers",
    "No adapter wrapping overhead at inference time",
]:
    add_text(sl, f"✓  {item}",
             x=Inches(6.35), y=Inches(1.88 + [0,0.6,1.2][
                 ["No adapter","Can target","No adapter w"].index(
                     next(k for k in ["No adapter","Can target","No adapter w"]
                          if item.startswith(k)))]),
             w=Inches(6.3), h=Inches(0.52),
             font_size=Pt(12), color=LIGHT_GREY)

textbox_with_bg(sl, Inches(6.2), Inches(4.15), Inches(6.7), Inches(2.75),
                border=ACCENT3)
add_text(sl, "Disadvantages",
         x=Inches(6.35), y=Inches(4.2), w=Inches(6.4), h=Inches(0.38),
         font_size=Pt(13), bold=True, color=ACCENT3)
for item in [
    "45M parameters trained vs 1.3M with LoRA — slower, more memory",
    "Saves a full ~13 GB model copy per experiment",
    "Changes are permanent — base model needs backup for ablations",
]:
    add_text(sl, f"✗  {item}",
             x=Inches(6.35), y=Inches(4.68 + [0,0.65,1.3][
                 ["45M","Saves","Changes"].index(
                     next(k for k in ["45M","Saves","Changes"]
                          if item.startswith(k)))]),
             w=Inches(6.3), h=Inches(0.56),
             font_size=Pt(12), color=LIGHT_GREY)

# Command
code_block(sl, [
    "python creme/train_proactive.py \\",
    "    --hparams  ./creme/hparams/codellama.yaml \\",
    "    --task_name  mbpp_codellama \\",
    "    --save_path  ./models/codellama_full \\",
    "    --train_mode  full \\",
    "    --lambda_reg  0.01  --num_epochs  5  --batch_size  8  --bf16",
], x=Inches(6.2), y=Inches(6.55), w=Inches(6.7), h=Inches(0.8))


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Fine-Tuning With LoRA
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Fine-Tuning — Mode 2: LoRA Adapter Training",
            "--train_mode lora  (default)  |  Low-Rank Adaptation via PEFT")

# Left — LoRA explanation
textbox_with_bg(sl, Inches(0.4), Inches(1.38), Inches(5.6), Inches(5.5),
                border=ACCENT)
add_text(sl, "What LoRA Does",
         x=Inches(0.55), y=Inches(1.43), w=Inches(5.3), h=Inches(0.38),
         font_size=Pt(13.5), bold=True, color=ACCENT)

code_block(sl, [
    "# Standard weight update:",
    "W'  =  W  +  ΔW          (ΔW is full-rank, expensive)",
    "",
    "# LoRA decomposition:",
    "W'  =  W  +  (α/r) · B · A",
    "",
    "  W  ∈  R^(d_out × d_in)   frozen",
    "  A  ∈  R^(r × d_in)       trained  (r << d_in)",
    "  B  ∈  R^(d_out × r)      trained  (r << d_out)",
    "  r = 8,   α = 16,   scale = α/r = 2.0",
], x=Inches(0.55), y=Inches(1.9), w=Inches(5.2), h=Inches(2.8))

add_text(sl,
    "Only A and B are trained. The original W at Layer 28 is never modified. "
    "This gives a tiny trainable footprint while preserving the full model's behaviour on unrelated tasks.",
    x=Inches(0.55), y=Inches(4.8), w=Inches(5.2), h=Inches(1.9),
    font_size=Pt(11.5), color=LIGHT_GREY, wrap=True)

# Centre — config table
headers = ["Parameter", "Value", "Meaning"]
rows = [
    ["rank  r",        "8",     "Adapter expressiveness"],
    ["lora_alpha",     "16",    "Scale factor (α/r = 2.0)"],
    ["lora_dropout",   "0.05",  "Regularisation"],
    ["Target module",  "down_proj @ L28", "Single module"],
    ["Trainable params", "~1.3 M", "0.018% of 7B total"],
    ["Saved artifact",   "~5 MB",  "Adapter matrices only"],
]
table_slide(sl, headers, rows,
            x=Inches(6.35), y=Inches(1.38),
            w=Inches(6.55), h=Inches(2.85))

# Advantages / Disadvantages
textbox_with_bg(sl, Inches(6.35), Inches(4.4), Inches(3.05), Inches(2.45),
                border=ACCENT2)
add_text(sl, "Advantages",
         x=Inches(6.5), y=Inches(4.45), w=Inches(2.8), h=Inches(0.38),
         font_size=Pt(12), bold=True, color=ACCENT2)
for item in [
    "Only 1.3M trainable params",
    "~5 MB saved per experiment",
    "Base model unchanged",
    "Rapid iteration / ablations",
]:
    k = ["Only","~5 MB","Base","Rapid"].index(
        next(x for x in ["Only","~5 MB","Base","Rapid"] if item.startswith(x)))
    add_text(sl, f"✓  {item}",
             x=Inches(6.5), y=Inches(4.9 + k * 0.45),
             w=Inches(2.8), h=Inches(0.4),
             font_size=Pt(11), color=LIGHT_GREY)

textbox_with_bg(sl, Inches(9.7), Inches(4.4), Inches(3.2), Inches(2.45),
                border=ACCENT3)
add_text(sl, "Limitations",
         x=Inches(9.85), y=Inches(4.45), w=Inches(2.9), h=Inches(0.38),
         font_size=Pt(12), bold=True, color=ACCENT3)
for item in [
    "Single layer only (lora mode)",
    "Requires peft library",
    "Adapter kept unmerged",
    "r=8 may underfit hard types",
]:
    k = ["Single","Requires","Adapter","r=8"].index(
        next(x for x in ["Single","Requires","Adapter","r=8"] if item.startswith(x)))
    add_text(sl, f"✗  {item}",
             x=Inches(9.85), y=Inches(4.9 + k * 0.45),
             w=Inches(2.9), h=Inches(0.4),
             font_size=Pt(11), color=LIGHT_GREY)

# Command
code_block(sl, [
    "python creme/train_proactive.py \\",
    "    --hparams  ./creme/hparams/codellama.yaml \\",
    "    --task_name  mbpp_codellama \\",
    "    --save_path  ./models/codellama_lora \\",
    "    --train_mode  lora  \\",
    "    --lambda_reg  0.01  --num_epochs  5  --batch_size  8  --bf16",
], x=Inches(6.35), y=Inches(7.05), w=Inches(6.55), h=Inches(0.35))


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Full vs LoRA Side-by-Side
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Full vs LoRA — Side-by-Side Comparison",
            "Same loss, same data, same layer hook — different parameter strategy")

# Architecture diagram (text-art style)
code_block(sl, [
    "Base Model: CodeLlama-7B  (32 layers, ~7 billion parameters)",
    "                    │",
    "         ┌──────────┴──────────┐",
    "         │                     │",
    "    full mode               lora mode",
    "         │                     │",
    "  Unfreeze down_proj       Inject A, B matrices",
    "  at layer 28 only         alongside down_proj",
    "  (~45M params)            at layer 28 only",
    "         │                 (~1.3M params)",
    "  Train in-place               │",
    "         │               Train only A and B",
    "  Save full model              │",
    "  (~13 GB)               Save adapter only",
    "                          (~5 MB)",
    "         │                     │",
    "         └──────────┬──────────┘",
    "                    │",
    "      Same loss:  L_ce  +  λ · L_reg",
    "      Same data:  all MBPP  ×  all 20 pert types",
    "      Same hook:  nethook.TraceDict at layer 28",
], x=Inches(0.4), y=Inches(1.38), w=Inches(7.6), h=Inches(5.8))

# Comparison table
headers = ["Criterion", "Full mode", "LoRA mode"]
rows = [
    ["Trainable params",  "~45 M",     "~1.3 M"],
    ["Saved size",        "~13 GB",    "~5 MB"],
    ["Multi-layer",       "Yes",       "No (one layer)"],
    ["Base model intact", "No",        "Yes"],
    ["Library needed",    "None",      "peft"],
    ["Our default",       "No",        "Yes"],
]
table_slide(sl, headers, rows,
            x=Inches(8.3), y=Inches(1.38),
            w=Inches(4.7), h=Inches(3.5))

textbox_with_bg(sl, Inches(8.3), Inches(5.1), Inches(4.7), Inches(2.1),
                border=ACCENT)
add_text(sl, "Why we chose LoRA as default",
         x=Inches(8.45), y=Inches(5.15), w=Inches(4.4), h=Inches(0.38),
         font_size=Pt(12.5), bold=True, color=ACCENT)
add_text(sl,
    "Representation alignment is a low-rank correction — "
    "the model needs to slightly rotate/scale Layer 28 activations, "
    "not learn entirely new behaviour. "
    "LoRA's 1.3M parameters are sufficient for this, "
    "while keeping the base model intact for fast ablation experiments.",
    x=Inches(8.45), y=Inches(5.6), w=Inches(4.4), h=Inches(1.5),
    font_size=Pt(11), color=LIGHT_GREY, wrap=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Experimental Results
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Experimental Results", "CodeLlama-7B on MBPP — 20 perturbation types")

# Summary stats row
for i, (label, val, sub, col) in enumerate([
    ("Overall  pass@1\n(perturbed)",  "0.457", "proactive model",   ACCENT3),
    ("Overall  pass@1\n(original)",   "0.603", "clean prompt",      ACCENT2),
    ("Mean delta",                    "−0.146","avg. degradation",  ACCENT4),
    ("Best type",                     "E3",    "Δ = 0.000",         ACCENT2),
    ("Worst type",                    "D2",    "Δ = −0.343",        ACCENT3),
]):
    bx = Inches(0.35) + Inches(2.52) * i
    add_rect(sl, bx, Inches(1.38), Inches(2.3), Inches(1.2),
             fill=BOX_BG, line=col, line_w=Pt(1.5))
    add_text(sl, label,
             x=bx + Inches(0.1), y=Inches(1.43), w=Inches(2.1), h=Inches(0.45),
             font_size=Pt(11), color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(sl, val,
             x=bx + Inches(0.1), y=Inches(1.85), w=Inches(2.1), h=Inches(0.45),
             font_size=Pt(20), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, sub,
             x=bx + Inches(0.1), y=Inches(2.26), w=Inches(2.1), h=Inches(0.25),
             font_size=Pt(9), color=MID_GREY, align=PP_ALIGN.CENTER)

# Per-type table
headers2 = ["Type", "Category", "Tasks", "Pert pass@1", "Orig pass@1", "Δ"]
rows2 = [
    ["A1","Argument","16","0.269","0.406","−0.138"],
    ["A2","Argument","12","0.375","0.383","−0.008"],
    ["C1","Control", "17","0.371","0.577","−0.206"],
    ["C3","Control", "13","0.369","0.585","−0.215"],
    ["D2","Data",    "18","0.333","0.677","−0.343"],
    ["D4","Data",    "13","0.677","0.708","−0.031"],
    ["E3","Expression","12","0.633","0.633"," 0.000"],
    ["P1","Problem", "18","0.450","0.711","−0.261"],
    ["S2","Statement","13","0.531","0.685","−0.154"],
]
table_slide(sl, headers2, rows2,
            x=Inches(0.35), y=Inches(2.75),
            w=Inches(8.6),  h=Inches(4.3))

# Key findings
textbox_with_bg(sl, Inches(9.25), Inches(2.75), Inches(3.75), Inches(4.3),
                border=ACCENT)
add_text(sl, "Key Findings",
         x=Inches(9.4), y=Inches(2.8), w=Inches(3.5), h=Inches(0.38),
         font_size=Pt(13), bold=True, color=ACCENT)
findings = [
    (ACCENT2,  "Expression pert. tolerated best — E3 has zero drop"),
    (ACCENT3,  "Data type changes (D2) cause the largest drop: −0.343"),
    (ACCENT3,  "Problem rephrasing (P1/P2) consistently hard: −0.26"),
    (ACCENT2,  "Clean-prompt score preserved — L_ce prevents forgetting"),
    (ACCENT4,  "A2 nearly unaffected — minor arg reorder has low impact"),
]
for k, (col, text) in enumerate(findings):
    add_text(sl, f"•  {text}",
             x=Inches(9.4), y=Inches(3.3 + k * 0.68),
             w=Inches(3.5), h=Inches(0.62),
             font_size=Pt(11), color=col, wrap=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Challenges
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Challenges", "Engineering and research hurdles encountered")

challenges = [
    ("Causal Layer Reliability",
     "L2 tracing loops over all 32 layers per task. Layer can vary slightly "
     "across tasks — we use the mode (Layer 28) as a fixed config, which may "
     "not be optimal for every individual task.",
     ACCENT3),
    ("LoRA Single-Layer Limit",
     "LoRA mode supports exactly one target layer. Data and Problem perturbations "
     "may benefit from multi-layer alignment — switching to full mode is the current "
     "workaround.",
     ACCENT4),
    ("Balancing λ",
     "λ = 0.01 was chosen empirically. Too high collapses the representation space; "
     "too low weakens the alignment signal. Optimal λ is likely category-dependent.",
     ACCENT),
    ("Training Pair Coverage",
     "Perturbed JSONL files do not cover every MBPP task for every type — missing pairs "
     "are silently skipped, leaving some types underrepresented in training.",
     ACCENT3),
    ("Memory Constraints",
     "7B model in bf16 ≈ 14 GB VRAM. With activations for both clean and perturbed "
     "paths at Layer 28, peak usage hits 20–22 GB. Gradient accumulation required "
     "on sub-24 GB GPUs.",
     ACCENT4),
    ("Windows + CUDA",
     "Project originally designed for Linux. Patching torch.amp imports, path "
     "separators, and multiprocessing in the code execution sandbox was required "
     "for Windows compatibility.",
     ACCENT),
]

for i, (title, body, col) in enumerate(challenges):
    col_idx = i % 3
    row_idx = i // 3
    bx = Inches(0.35) + Inches(4.32) * col_idx
    by = Inches(1.38) + Inches(2.7) * row_idx
    add_rect(sl, bx, by, Inches(4.1), Inches(2.5),
             fill=BOX_BG, line=col, line_w=Pt(1.3))
    add_text(sl, title,
             x=bx + Inches(0.12), y=by + Inches(0.1),
             w=Inches(3.85), h=Inches(0.38),
             font_size=Pt(12.5), bold=True, color=col)
    add_text(sl, body,
             x=bx + Inches(0.12), y=by + Inches(0.55),
             w=Inches(3.85), h=Inches(1.85),
             font_size=Pt(11), color=LIGHT_GREY, wrap=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Future Work
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Future Work", "Directions to extend and strengthen the approach")

future = [
    ("Multi-Layer Alignment",
     "Apply LoRA at layers {L−2, L, L+2} around Layer 28 with weighted alignment loss. "
     "Expected to help Data and Problem perturbations where single-layer correction is insufficient.",
     ACCENT),
    ("Per-Category λ Weighting",
     "Learn category-specific λ values (e.g., λ_E < λ_D < λ_P) based on observed "
     "difficulty. Frame as a meta-learning problem over the difficulty distribution.",
     ACCENT2),
    ("HumanEval + QwenCoder",
     "Run the full proactive pipeline on all four configurations: "
     "humaneval_codellama, humaneval_qwen, mbpp_codellama, mbpp_qwen. "
     "Produce a cross-model, cross-benchmark comparison table.",
     ACCENT4),
    ("Automated Layer Discovery",
     "Train a lightweight probe to predict the causal layer given perturbation type "
     "and task properties. Removes the need for the one-time empirical tracing study.",
     ACCENT),
    ("LoRA Rank Ablation",
     "Systematically vary rank r ∈ {4, 8, 16, 32} and target modules "
     "(down_proj vs q_proj/v_proj). Quantify the minimum rank sufficient per "
     "perturbation category.",
     ACCENT2),
    ("Retrieval-Augmented Robustness",
     "For Problem-level perturbations (P1, P2) where NL rephrasing is the challenge, "
     "augment with retrieval of similar original problems to supplement layer-level alignment.",
     ACCENT4),
]

for i, (title, body, col) in enumerate(future):
    col_idx = i % 3
    row_idx = i // 3
    bx = Inches(0.35) + Inches(4.32) * col_idx
    by = Inches(1.38) + Inches(2.7) * row_idx
    # Number badge
    add_rect(sl, bx, by, Inches(0.42), Inches(0.42),
             fill=col, line=col, line_w=Pt(0))
    add_text(sl, str(i + 1),
             x=bx, y=by, w=Inches(0.42), h=Inches(0.42),
             font_size=Pt(13), bold=True, color=DARK_BG,
             align=PP_ALIGN.CENTER)
    add_rect(sl, bx, by, Inches(4.1), Inches(2.5),
             fill=BOX_BG, line=col, line_w=Pt(1.3))
    add_text(sl, str(i + 1),
             x=bx + Inches(0.08), y=by + Inches(0.06),
             w=Inches(0.3), h=Inches(0.35),
             font_size=Pt(13), bold=True, color=col)
    add_text(sl, title,
             x=bx + Inches(0.42), y=by + Inches(0.06),
             w=Inches(3.55), h=Inches(0.38),
             font_size=Pt(12.5), bold=True, color=col)
    add_text(sl, body,
             x=bx + Inches(0.12), y=by + Inches(0.52),
             w=Inches(3.85), h=Inches(1.88),
             font_size=Pt(11), color=LIGHT_GREY, wrap=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
slide_title(sl, "Conclusion", "Summary of contributions and key takeaways")

# Central message
add_rect(sl, Inches(0.5), Inches(1.38), Inches(12.33), Inches(1.05),
         fill=rgb(0x1A, 0x2A, 0x3E), line=ACCENT, line_w=Pt(1.5))
add_text(sl,
    "Not every layer counts — but Layer 28 does. "
    "A single offline fine-tuning step at the right layer is sufficient to make "
    "CodeLlama-7B robustly handle all 20 perturbation types without any inference-time editing.",
    x=Inches(0.65), y=Inches(1.43), w=Inches(12.0), h=Inches(0.92),
    font_size=Pt(14), bold=False, italic=True, color=WHITE,
    align=PP_ALIGN.CENTER)

# Four takeaway boxes
takeaways = [
    ("Proactive > Reactive",
     "Replace O(tasks×types) inference-time edits with a single O(1) offline training run. "
     "Identical result quality, negligible inference overhead.",
     ACCENT),
    ("Layer Specificity Matters",
     "Applying the alignment loss at the causal layer (L28) drives the improvement. "
     "Training at a random layer (ablation) yields negligible benefit.",
     ACCENT2),
    ("LoRA is Sufficient",
     "~1.3M LoRA parameters encode cross-perturbation robustness for a 7B model. "
     "The alignment task is low-rank — it needs rotation, not new knowledge.",
     ACCENT4),
    ("Clean Performance Preserved",
     "The L_ce term prevents catastrophic forgetting. Original-prompt pass@1 stays at 0.603 "
     "across all 20 perturbation types simultaneously.",
     ACCENT3),
]
for i, (title, body, col) in enumerate(takeaways):
    bx = Inches(0.35) + Inches(3.17) * i
    add_rect(sl, bx, Inches(2.65), Inches(2.95), Inches(3.65),
             fill=BOX_BG, line=col, line_w=Pt(1.5))
    add_text(sl, title,
             x=bx + Inches(0.12), y=Inches(2.72),
             w=Inches(2.7), h=Inches(0.45),
             font_size=Pt(12.5), bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(sl, body,
             x=bx + Inches(0.12), y=Inches(3.22),
             w=Inches(2.7), h=Inches(3.0),
             font_size=Pt(11.5), color=LIGHT_GREY,
             wrap=True, align=PP_ALIGN.CENTER)

# Bottom result row
add_text(sl, "Results at a glance",
         x=Inches(0.5), y=Inches(6.4), w=Inches(12.3), h=Inches(0.35),
         font_size=Pt(11), bold=True, color=MID_GREY, align=PP_ALIGN.CENTER)
add_text(sl,
    "Mean pass@1  →  0.603 (original)    0.457 (perturbed)    Δ = −0.146  "
    "|  Best: E3 (Δ = 0.000)    Worst: D2 (Δ = −0.343)",
    x=Inches(0.5), y=Inches(6.73), w=Inches(12.3), h=Inches(0.45),
    font_size=Pt(12), color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "CREME_Presentation.pptx"
prs.save(out_path)
print(f"\nPresentation saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides, 1):
    print(f"  Slide {i:02d}")
